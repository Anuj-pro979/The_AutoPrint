# wo3_autoprint_streamlit_firestore_sender_fixed.py
# Cloud-compatible Streamlit sender that uploads chunked base64 docs + manifest to Firestore
# Optimized for Streamlit Cloud with proper error handling and fallbacks
#
# Run: streamlit run wo3_autoprint_streamlit_firestore_sender_fixed.py

import streamlit as st
import streamlit.components.v1 as components
import os
import tempfile
import base64
import time
import json
import logging
import traceback
import subprocess
import platform
from typing import Optional, List, Dict, Any
from dataclasses import dataclass
from fpdf import FPDF
from PIL import Image
from pathlib import Path
import hashlib
import datetime
import uuid
import webbrowser
import io
import zipfile
import xml.etree.ElementTree as ET

# Firestore
try:
    import firebase_admin
    from firebase_admin import credentials, firestore
    FIRESTORE_AVAILABLE = True
except ImportError:
    firebase_admin = None
    credentials = None
    firestore = None
    FIRESTORE_AVAILABLE = False

# PDF processing - use modern pypdf instead of deprecated PyPDF2
try:
    from pypdf import PdfReader
    PDF_READER_AVAILABLE = True
except ImportError:
    try:
        from PyPDF2 import PdfReader
        PDF_READER_AVAILABLE = True
    except ImportError:
        PdfReader = None
        PDF_READER_AVAILABLE = False

# QR generation
try:
    import qrcode
    QR_AVAILABLE = True
except ImportError:
    QR_AVAILABLE = False

# python-docx for DOCX text extraction
try:
    import docx as python_docx
    PYTHON_DOCX_AVAILABLE = True
except ImportError:
    python_docx = None
    PYTHON_DOCX_AVAILABLE = False

# python-pptx for PPTX text extraction
try:
    from pptx import Presentation as PptxPresentation
    PYTHON_PPTX_AVAILABLE = True
except ImportError:
    PptxPresentation = None
    PYTHON_PPTX_AVAILABLE = False

# --------- Improved Logging ----------
def setup_logger():
    logger = logging.getLogger("autoprint_sender")
    if not logger.handlers:
        logger.setLevel(logging.DEBUG)
        handler = logging.StreamHandler()
        formatter = logging.Formatter("%(asctime)s - %(levelname)s - %(message)s")
        handler.setFormatter(formatter)
        logger.addHandler(handler)
    return logger

logger = setup_logger()

def log_message(msg: str, level: str = "info"):
    """Improved logging with Streamlit integration"""
    if level == "debug":
        logger.debug(msg)
    elif level == "warning":
        logger.warning(msg)
        st.warning(f"⚠️ {msg}")
    elif level == "error":
        logger.error(msg)
        st.error(f"❌ {msg}")
    else:
        logger.info(msg)
        st.info(f"ℹ️ {msg}")

# --------- Utilities ----------
def safe_remove(path: str):
    """Safely remove a file with error handling"""
    try:
        if path and os.path.exists(path):
            os.unlink(path)
    except Exception as e:
        logger.warning(f"Failed to remove {path}: {e}")

def retry_with_backoff(func, attempts=3, initial_delay=0.5, factor=2.0, *args, **kwargs):
    """Retry function with exponential backoff"""
    delay = initial_delay
    last_exc = None
    for i in range(attempts):
        try:
            return func(*args, **kwargs)
        except Exception as e:
            last_exc = e
            func_name = getattr(func, "__name__", str(func))
            logger.warning(f"Attempt {i+1}/{attempts} failed for {func_name}: {e}")
            if i < attempts - 1:  # Don't sleep on last attempt
                time.sleep(delay)
                delay *= factor
    
    func_name = getattr(func, "__name__", str(func))
    logger.error(f"All {attempts} attempts failed for {func_name}")
    if last_exc:
        raise last_exc
    return None

# --------- Data classes ----------
@dataclass
class PrintSettings:
    copies: int = 1
    color_mode: str = "Color"
    duplex: str = "Single-sided"
    paper_size: str = "A4"
    orientation: str = "Portrait"
    quality: str = "High"
    collate: bool = True
    staple: bool = False

@dataclass
class ConvertedFile:
    orig_name: str
    pdf_name: str
    pdf_bytes: bytes
    settings: PrintSettings
    original_bytes: Optional[bytes] = None
    conversion_method: str = "unknown"
    pages: int = 1

# --------- Improved FileConverter ----------
class FileConverter:
    SUPPORTED_TEXT_EXTENSIONS = {'.txt', '.md', '.rtf', '.html', '.htm', '.csv', '.log'}
    SUPPORTED_IMAGE_EXTENSIONS = {'.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.tif', '.webp', '.gif'}
    
    @staticmethod
    def create_text_pdf(text: str, title: str = "Document") -> bytes:
        """Create PDF from text with better formatting"""
        try:
            pdf = FPDF(unit='mm', format='A4')
            pdf.set_auto_page_break(auto=True, margin=15)
            pdf.add_page()
            
            # Add title
            pdf.set_font("Arial", 'B', 16)
            pdf.cell(0, 10, title, ln=True, align='C')
            pdf.ln(5)
            
            # Add content
            pdf.set_font("Arial", size=11)
            
            # Split text into paragraphs
            paragraphs = text.split('\n')
            
            for paragraph in paragraphs:
                if not paragraph.strip():
                    pdf.ln(3)
                    continue
                
                # Handle long paragraphs by splitting them
                max_chars_per_line = 85
                if len(paragraph) <= max_chars_per_line:
                    try:
                        pdf.multi_cell(0, 5, paragraph.encode('latin-1', 'replace').decode('latin-1'))
                    except:
                        pdf.multi_cell(0, 5, paragraph.encode('utf-8', 'replace').decode('utf-8', 'replace'))
                else:
                    # Split long paragraphs
                    words = paragraph.split(' ')
                    current_line = ""
                    
                    for word in words:
                        if len(current_line + word) <= max_chars_per_line:
                            current_line += word + " "
                        else:
                            if current_line:
                                try:
                                    pdf.multi_cell(0, 5, current_line.strip().encode('latin-1', 'replace').decode('latin-1'))
                                except:
                                    pdf.multi_cell(0, 5, current_line.strip())
                            current_line = word + " "
                    
                    if current_line:
                        try:
                            pdf.multi_cell(0, 5, current_line.strip().encode('latin-1', 'replace').decode('latin-1'))
                        except:
                            pdf.multi_cell(0, 5, current_line.strip())
                
                pdf.ln(2)
            
            return pdf.output(dest='S').encode('latin-1', errors='replace')
            
        except Exception as e:
            logger.error(f"Failed to create text PDF: {e}")
            # Fallback: create simple PDF with error message
            pdf = FPDF()
            pdf.add_page()
            pdf.set_font("Arial", size=12)
            pdf.multi_cell(0, 10, f"Error creating PDF from text: {str(e)}")
            return pdf.output(dest='S').encode('latin-1', errors='replace')

    @classmethod
    def convert_text_file(cls, file_content: bytes, filename: str) -> Optional[bytes]:
        """Convert text-based files to PDF"""
        try:
            # Try different encodings
            text = None
            for encoding in ['utf-8', 'utf-8-sig', 'latin-1', 'cp1252', 'ascii']:
                try:
                    text = file_content.decode(encoding, errors='ignore')
                    break
                except:
                    continue
            
            if not text:
                text = file_content.decode('utf-8', errors='replace')
            
            if not text.strip():
                text = f"Empty or unreadable file: {filename}"
            
            return cls.create_text_pdf(text, os.path.splitext(filename)[0])
            
        except Exception as e:
            logger.error(f"Text file conversion failed for {filename}: {e}")
            return cls.create_text_pdf(f"Error reading file {filename}: {str(e)}", filename)

    @classmethod
    def convert_image_file(cls, file_content: bytes, filename: str) -> Optional[bytes]:
        """Convert image files to PDF with better error handling"""
        try:
            with Image.open(io.BytesIO(file_content)) as img:
                # Handle different image modes
                if img.mode not in ('RGB', 'L'):
                    img = img.convert('RGB')
                
                # Resize if too large (memory optimization)
                max_dimension = 2000
                if img.width > max_dimension or img.height > max_dimension:
                    img.thumbnail((max_dimension, max_dimension), Image.Resampling.LANCZOS)
                
                # Create PDF
                pdf_buffer = io.BytesIO()
                img.save(pdf_buffer, format='PDF', optimize=True, quality=85)
                return pdf_buffer.getvalue()
                
        except Exception as e:
            logger.error(f"Image conversion failed for {filename}: {e}")
            # Create fallback PDF with error message
            return cls.create_text_pdf(f"Failed to convert image: {filename}\nError: {str(e)}", filename)

    @classmethod
    def extract_docx_text_xml(cls, file_content: bytes) -> str:
        """Extract text from DOCX using XML parsing (no external dependencies)"""
        try:
            with zipfile.ZipFile(io.BytesIO(file_content), 'r') as docx_zip:
                if 'word/document.xml' not in docx_zip.namelist():
                    return "No document.xml found in DOCX file"
                
                xml_content = docx_zip.read('word/document.xml')
                root = ET.fromstring(xml_content)
                
                # Extract text from all text nodes
                texts = []
                for elem in root.iter():
                    if elem.text and elem.text.strip():
                        # Skip namespace prefixes in tag names
                        tag_name = elem.tag.split('}')[-1] if '}' in elem.tag else elem.tag
                        if tag_name == 't':  # text elements
                            texts.append(elem.text)
                        elif tag_name in ['p', 'tc']:  # paragraphs and table cells
                            if elem.text:
                                texts.append('\n' + elem.text)
                
                return '\n'.join(texts) if texts else "No text content found in DOCX"
                
        except Exception as e:
            logger.error(f"DOCX XML extraction failed: {e}")
            return f"Error extracting DOCX content: {str(e)}"

    @classmethod
    def extract_docx_text_library(cls, file_content: bytes) -> str:
        """Extract text from DOCX using python-docx library"""
        try:
            if not PYTHON_DOCX_AVAILABLE:
                return None
            
            doc = python_docx.Document(io.BytesIO(file_content))
            paragraphs = []
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    paragraphs.append(paragraph.text)
            
            # Also extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            paragraphs.append(cell.text)
            
            return '\n\n'.join(paragraphs) if paragraphs else "No text content found"
            
        except Exception as e:
            logger.error(f"python-docx extraction failed: {e}")
            return None

    @classmethod
    def convert_docx_file(cls, file_content: bytes, filename: str) -> Optional[bytes]:
        """Convert DOCX files to PDF using multiple fallback methods"""
        try:
            # Method 1: Try python-docx library if available
            text = cls.extract_docx_text_library(file_content)
            
            # Method 2: Fallback to XML parsing
            if not text:
                text = cls.extract_docx_text_xml(file_content)
            
            if text and text.strip():
                return cls.create_text_pdf(text, os.path.splitext(filename)[0])
            else:
                # Final fallback
                return cls.create_text_pdf(
                    f"Unable to extract readable content from: {filename}\n\n"
                    "This DOCX file may contain complex formatting, images, or be corrupted.\n"
                    "Please try converting it to PDF manually or use a simpler format.",
                    filename
                )
            
        except Exception as e:
            logger.error(f"DOCX conversion failed for {filename}: {e}")
            return cls.create_text_pdf(f"Error processing DOCX file: {filename}\nError: {str(e)}", filename)

    @classmethod
    def convert_pptx_file(cls, file_content: bytes, filename: str) -> Optional[bytes]:
        """Convert PPTX files to PDF"""
        try:
            if not PYTHON_PPTX_AVAILABLE:
                return cls.create_text_pdf(
                    f"PPTX conversion not available for: {filename}\n\n"
                    "The python-pptx library is not installed.\n"
                    "Text content cannot be extracted from PowerPoint files.",
                    filename
                )
            
            prs = PptxPresentation(io.BytesIO(file_content))
            slides_text = []
            
            for i, slide in enumerate(prs.slides, 1):
                slide_content = [f"=== Slide {i} ==="]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text)
                
                if len(slide_content) == 1:  # Only the slide header
                    slide_content.append("[No text content on this slide]")
                
                slides_text.append('\n'.join(slide_content))
            
            if slides_text:
                full_text = '\n\n'.join(slides_text)
                return cls.create_text_pdf(full_text, os.path.splitext(filename)[0])
            else:
                return cls.create_text_pdf(f"No content found in presentation: {filename}", filename)
            
        except Exception as e:
            logger.error(f"PPTX conversion failed for {filename}: {e}")
            return cls.create_text_pdf(f"Error processing PPTX file: {filename}\nError: {str(e)}", filename)

    @classmethod
    def convert_uploaded_file_to_pdf(cls, uploaded_file) -> Optional[ConvertedFile]:
        """Main conversion method with comprehensive fallbacks"""
        if not uploaded_file:
            return None
        
        filename = uploaded_file.name
        file_content = uploaded_file.getvalue()
        suffix = os.path.splitext(filename)[1].lower()
        
        try:
            # Handle PDF files (pass through)
            if suffix == ".pdf":
                pages = count_pdf_pages(file_content)
                return ConvertedFile(
                    orig_name=filename,
                    pdf_name=filename,
                    pdf_bytes=file_content,
                    settings=PrintSettings(),
                    original_bytes=file_content,
                    conversion_method="passthrough",
                    pages=pages
                )
            
            pdf_bytes = None
            conversion_method = "unknown"
            
            # Handle text files
            if suffix in cls.SUPPORTED_TEXT_EXTENSIONS:
                pdf_bytes = cls.convert_text_file(file_content, filename)
                conversion_method = "text"
            
            # Handle image files
            elif suffix in cls.SUPPORTED_IMAGE_EXTENSIONS:
                pdf_bytes = cls.convert_image_file(file_content, filename)
                conversion_method = "image"
            
            # Handle DOCX files
            elif suffix == ".docx":
                pdf_bytes = cls.convert_docx_file(file_content, filename)
                conversion_method = "docx"
            
            # Handle PPTX files
            elif suffix == ".pptx":
                pdf_bytes = cls.convert_pptx_file(file_content, filename)
                conversion_method = "pptx"
            
            # Unsupported format
            else:
                pdf_bytes = cls.create_text_pdf(
                    f"Unsupported file format: {suffix}\n\n"
                    f"File: {filename}\n"
                    f"Size: {len(file_content)} bytes\n\n"
                    "Supported formats:\n"
                    "• PDF (passthrough)\n"
                    "• Text: .txt, .md, .rtf, .html, .htm, .csv, .log\n"
                    "• Images: .png, .jpg, .jpeg, .bmp, .tiff, .webp, .gif\n"
                    "• Documents: .docx, .pptx\n\n"
                    "Please convert your file to a supported format.",
                    filename
                )
                conversion_method = "unsupported"
            
            if pdf_bytes:
                pages = count_pdf_pages(pdf_bytes)
                pdf_name = os.path.splitext(filename)[0] + ".pdf"
                
                return ConvertedFile(
                    orig_name=filename,
                    pdf_name=pdf_name,
                    pdf_bytes=pdf_bytes,
                    settings=PrintSettings(),
                    original_bytes=file_content,
                    conversion_method=conversion_method,
                    pages=pages
                )
            
            return None
            
        except Exception as e:
            logger.error(f"File conversion failed for {filename}: {e}")
            # Create error PDF as final fallback
            try:
                error_pdf = cls.create_text_pdf(
                    f"Conversion Error\n\n"
                    f"File: {filename}\n"
                    f"Error: {str(e)}\n\n"
                    "Please try a different file format or contact support.",
                    filename
                )
                return ConvertedFile(
                    orig_name=filename,
                    pdf_name=f"ERROR_{filename}.pdf",
                    pdf_bytes=error_pdf,
                    settings=PrintSettings(),
                    original_bytes=file_content,
                    conversion_method="error",
                    pages=1
                )
            except:
                return None

# --------- PDF Page Counting ----------
def count_pdf_pages(pdf_bytes: Optional[bytes]) -> int:
    """Count pages in PDF with better error handling"""
    if not pdf_bytes:
        return 1
    
    if not PDF_READER_AVAILABLE:
        logger.warning("PDF reader not available, defaulting to 1 page")
        return 1
    
    try:
        reader = PdfReader(io.BytesIO(pdf_bytes))
        return len(reader.pages)
    except Exception as e:
        logger.warning(f"Failed to count PDF pages: {e}")
        # Try to estimate based on file size (rough estimate)
        size_kb = len(pdf_bytes) / 1024
        if size_kb < 50:
            return 1
        elif size_kb < 200:
            return max(1, int(size_kb / 50))
        else:
            return max(1, int(size_kb / 100))

# --------- Streamlit Configuration ----------
st.set_page_config(
    page_title="Autoprint (Firestore)", 
    layout="wide", 
    page_icon="🖨️",
    initial_sidebar_state="expanded"
)

# Custom CSS for better appearance
st.markdown("""
<style>
    .main {
        padding-top: 1rem;
    }
    .stAlert {
        margin: 0.5rem 0;
    }
    .stButton > button {
        width: 100%;
    }
    .upload-section {
        border: 2px dashed #ccc;
        border-radius: 10px;
        padding: 1rem;
        margin: 1rem 0;
        text-align: center;
    }
    .file-item {
        border: 1px solid #e0e0e0;
        border-radius: 5px;
        padding: 0.5rem;
        margin: 0.5rem 0;
    }
    .success-message {
        background-color: #d4edda;
        border: 1px solid #c3e6cb;
        color: #155724;
        padding: 0.75rem 1.25rem;
        margin-bottom: 1rem;
        border-radius: 0.25rem;
    }
</style>
""", unsafe_allow_html=True)

# --------- Header ----------
st.markdown("<h1 style='text-align:center; margin-bottom:2rem;'>🖨️ Autoprint - Firestore Sender</h1>", unsafe_allow_html=True)

# --------- Initialize Session State ----------
def init_session_state():
    defaults = {
        'converted_files': [],
        'payinfo': None,
        'status': "",
        'process_complete': False,
        'user_name': "",
        'user_id': str(uuid.uuid4())[:8],
        'pricing': {
            "price_bw_per_page": 2.00,
            "price_color_per_page": 5.00,
            "price_duplex_discount": 0.9,
            "min_charge": 5.00,
            "currency": "INR",
            "owner_upi": "owner@upi"
        }
    }
    
    for key, value in defaults.items():
        if key not in st.session_state:
            st.session_state[key] = value

init_session_state()

def set_status(message: str):
    """Update status with timestamp"""
    timestamp = datetime.datetime.now().strftime('%H:%M:%S')
    st.session_state.status = f"{timestamp} - {message}"

# --------- Firestore Initialization ----------
COLLECTION = "files"
CHUNK_SIZE = 200_000  # characters per chunk

db = None
FIRESTORE_OK = False
FIRESTORE_ERR = None

def init_firestore():
    global db, FIRESTORE_OK, FIRESTORE_ERR
    
    if not FIRESTORE_AVAILABLE:
        FIRESTORE_ERR = "firebase_admin package not installed"
        return
    
    try:
        # Get service account from Streamlit secrets
        if not hasattr(st, "secrets") or "firebase_service_account" not in st.secrets:
            raise RuntimeError("Add 'firebase_service_account' to Streamlit Secrets")
        
        service_account_info = st.secrets["firebase_service_account"]
        
        # Handle both dict and JSON string formats
        if isinstance(service_account_info, str):
            service_account_info = json.loads(service_account_info)
        
        # Fix newlines in private key
        if "private_key" in service_account_info:
            service_account_info["private_key"] = service_account_info["private_key"].replace("\\n", "\n")
        
        # Initialize Firebase app if not already done
        try:
            app = firebase_admin.get_app()
        except ValueError:
            cred = credentials.Certificate(service_account_info)
            app = firebase_admin.initialize_app(cred)
        
        db = firestore.client()
        FIRESTORE_OK = True
        set_status("Firestore initialized successfully")
        
    except Exception as e:
        FIRESTORE_OK = False
        FIRESTORE_ERR = str(e)
        set_status(f"Firestore initialization failed: {e}")
        logger.error(f"Firestore init error: {e}")

# Initialize Firestore
init_firestore()

# --------- Utility Functions ----------
def sha256_bytes(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()

def meta_doc_id(file_id: str) -> str:
    return f"{file_id}_meta"

def chunk_doc_id(file_id: str, chunk_index: int) -> str:
    return f"{file_id}_{chunk_index}"

def calculate_amount(pricing_config: dict, pages: int, copies: int = 1, color: bool = False, duplex: bool = False) -> float:
    try:
        price_per_page = pricing_config.get("price_color_per_page", 5.0) if color else pricing_config.get("price_bw_per_page", 2.0)
        amount = pages * price_per_page * copies
        
        if duplex:
            amount *= pricing_config.get("price_duplex_discount", 0.9)
        
        min_charge = pricing_config.get("min_charge", 5.0)
        return max(amount, min_charge)
        
    except Exception as e:
        logger.error(f"Amount calculation failed: {e}")
        return pricing_config.get("min_charge", 5.0)

def generate_upi_uri(upi_id: str, amount: float, note: str = None) -> str:
    params = [f"pa={upi_id}", f"am={amount:.2f}"]
    if note:
        params.append(f"tn={note}")
    return "upi://pay?" + "&".join(params)

# --------- File Upload and Processing ----------
def upload_files_to_firestore(converted_files: List[ConvertedFile], job_settings: dict):
    """Upload files to Firestore with progress tracking"""
    
    if not FIRESTORE_OK or not db:
        st.error("❌ Firestore not available. Cannot upload files.")
        return False
    
    if not converted_files:
        st.error("❌ No files to upload.")
        return False
    
    try:
        job_id = str(uuid.uuid4())[:12]
        set_status(f"Starting upload for job {job_id}")
        
        # Prepare file metadata
        files_metadata = []
        total_chunks = 0
        
        for cf in converted_files:
            file_id = str(uuid.uuid4())[:8]
            pdf_data = cf.pdf_bytes
            
            if not pdf_data:
                st.warning(f"⚠️ No PDF data for {cf.orig_name}, skipping")
                continue
            
            # Convert to base64 and chunk
            b64_data = base64.b64encode(pdf_data).decode('utf-8')
            chunks = [b64_data[i:i+CHUNK_SIZE] for i in range(0, len(b64_data), CHUNK_SIZE)]
            
            file_meta = {
                "file_id": file_id,
                "filename": cf.pdf_name,
                "orig_filename": cf.orig_name,
                "size_bytes": len(pdf_data),
                "pages": cf.pages,
                "conversion_method": cf.conversion_method,
                "settings": {
                    "copies": job_settings.get("copies", 1),
                    "color_mode": job_settings.get("color_mode", "Color"),
                    "duplex": cf.settings.duplex,
                    "paper_size": cf.settings.paper_size,
                    "orientation": cf.settings.orientation,
                    "collate": cf.settings.collate
                },
                "chunks": chunks,
                "total_chunks": len(chunks),
                "sha256": sha256_bytes(pdf_data),
                "job_id": job_id
            }
            
            files_metadata.append(file_meta)
            total_chunks += len(chunks)
        
        if not files_metadata:
            st.error("❌ No valid files to upload after processing.")
            return False
        
        # Upload progress bar
        progress_bar = st.progress(0.0)
        status_text = st.empty()
        uploaded_chunks = 0
        
        # Upload chunks for each file
        for file_meta in files_metadata:
            file_id = file_meta["file_id"]
            filename = file_meta["filename"]
            
            status_text.text(f"Uploading {filename}...")
            
            # Upload file chunks
            for chunk_index, chunk_data in enumerate(file_meta["chunks"]):
                chunk_doc_id_str = chunk_doc_id(file_id, chunk_index)
                
                def upload_chunk():
                    doc_ref = db.collection(COLLECTION).document(chunk_doc_id_str)
                    doc_ref.set({
                        "data": chunk_data,
                        "chunk_index": chunk_index,
                        "file_id": file_id,
                        "timestamp": datetime.datetime.now()
                    })
                
                retry_with_backoff(upload_chunk, attempts=3)
                uploaded_chunks += 1
                
                # Update progress
                progress = uploaded_chunks / total_chunks
                progress_bar.progress(progress)
            
            # Upload file metadata
            meta_doc = {
                "total_chunks": file_meta["total_chunks"],
                "file_name": file_meta["filename"],
                "orig_filename": file_meta["orig_filename"],
                "sha256": file_meta["sha256"],
                "file_size_bytes": file_meta["size_bytes"],
                "pages": file_meta["pages"],
                "conversion_method": file_meta["conversion_method"],
                "settings": file_meta["settings"],
                "user_name": st.session_state.get("user_name", ""),
                "user_id": st.session_state.get("user_id", ""),
                "job_id": job_id,
                "timestamp": datetime.datetime.now(),
                "status": "uploaded"
            }
            
            def upload_metadata():
                meta_doc_ref = db.collection(COLLECTION).document(meta_doc_id(file_id))
                meta_doc_ref.set(meta_doc, merge=True)
            
            retry_with_backoff(upload_metadata, attempts=3)
            
            set_status(f"Uploaded {filename} ({file_meta['total_chunks']} chunks)")
        
        progress_bar.progress(1.0)
        status_text.text("✅ Upload completed!")
        
        # Wait for payment info from receiver
        poll_for_payment_info(files_metadata, job_settings)
        
        return True
        
    except Exception as e:
        logger.error(f"Upload failed: {e}")
        st.error(f"❌ Upload failed: {str(e)}")
        return False

def poll_for_payment_info(files_metadata: List[dict], job_settings: dict):
    """Poll Firestore for payment information from receiver"""
    
    set_status("Waiting for payment information from receiver...")
    
    # Show local estimate first
    show_local_estimate(files_metadata, job_settings)
    
    # Poll for official payment info
    poll_start = time.time()
    max_poll_time = 120  # 2 minutes
    
    progress_container = st.container()
    
    while time.time() - poll_start < max_poll_time:
        try:
            # Check each file's metadata for payment info
            for file_meta in files_metadata:
                file_id = file_meta["file_id"]
                
                doc_ref = db.collection(COLLECTION).document(meta_doc_id(file_id))
                doc_snapshot = doc_ref.get()
                
                if doc_snapshot.exists:
                    doc_data = doc_snapshot.to_dict()
                    
                    if doc_data.get("payinfo"):
                        st.session_state.payinfo = doc_data["payinfo"]
                        set_status("✅ Received official payment information!")
                        return
            
            # Update polling status
            elapsed = int(time.time() - poll_start)
            remaining = max_poll_time - elapsed
            
            with progress_container:
                st.info(f"⏳ Polling for payment info... ({remaining}s remaining)")
            
            time.sleep(2)
            
        except Exception as e:
            logger.error(f"Polling error: {e}")
            break
    
    # Timeout reached
    set_status("⚠️ Timeout waiting for official payment info. Using local estimate.")

def show_local_estimate(files_metadata: List[dict], job_settings: dict):
    """Show local payment estimate while waiting for official info"""
    
    try:
        pricing = st.session_state.get("pricing", {})
        total_amount = 0.0
        total_pages = 0
        
        for file_meta in files_metadata:
            pages = file_meta["pages"]
            copies = job_settings.get("copies", 1)
            
            # Determine if color printing
            color_mode = job_settings.get("color_mode", "Color")
            is_color = "color" in color_mode.lower()
            
            # Check for duplex
            duplex_setting = file_meta["settings"].get("duplex", "Single-sided")
            is_duplex = "duplex" in duplex_setting.lower() or "two" in duplex_setting.lower()
            
            file_amount = calculate_amount(pricing, pages, copies, is_color, is_duplex)
            total_amount += file_amount
            total_pages += pages * copies
        
        # Create payment info object
        job_id = files_metadata[0]["job_id"]
        file_name = files_metadata[0]["filename"] if len(files_metadata) == 1 else "Multiple Files"
        
        payment_estimate = {
            "order_id": job_id,
            "file_name": file_name,
            "total_files": len(files_metadata),
            "pages": total_pages,
            "copies": job_settings.get("copies", 1),
            "amount": round(total_amount, 2),
            "currency": pricing.get("currency", "INR"),
            "owner_upi": pricing.get("owner_upi", "owner@upi"),
            "status": "estimated",
            "is_estimate": True
        }
        
        st.session_state.payinfo = payment_estimate
        set_status("💰 Local payment estimate ready")
        
    except Exception as e:
        logger.error(f"Failed to create payment estimate: {e}")

# --------- Payment Handling ----------
def handle_online_payment():
    """Handle UPI online payment"""
    payinfo = st.session_state.get("payinfo")
    if not payinfo:
        st.error("❌ No payment information available")
        return
    
    amount = payinfo.get("amount", 0)
    upi_id = payinfo.get("owner_upi")
    file_name = payinfo.get("file_name", "Print Job")
    
    if not upi_id:
        st.error("❌ UPI ID not available")
        return
    
    # Generate UPI URI
    upi_uri = generate_upi_uri(upi_id, amount, f"Print: {file_name}")
    
    st.success(f"💳 Please pay ₹{amount:.2f} via UPI")
    
    # Create payment link
    st.markdown(f"### [🚀 Open Payment App]({upi_uri})")
    
    # Generate QR code if available
    if QR_AVAILABLE:
        try:
            qr = qrcode.QRCode(version=1, box_size=8, border=2)
            qr.add_data(upi_uri)
            qr.make(fit=True)
            
            qr_img = qr.make_image(fill_color="black", back_color="white")
            
            # Convert PIL image to bytes for Streamlit
            img_buffer = io.BytesIO()
            qr_img.save(img_buffer, format='PNG')
            img_buffer.seek(0)
            
            st.image(img_buffer, width=250, caption="Scan with any UPI app")
            
        except Exception as e:
            logger.warning(f"QR code generation failed: {e}")
    
    # Try to open payment app automatically
    try:
        webbrowser.open(upi_uri)
    except:
        pass
    
    st.balloons()
    complete_payment_process()

def handle_offline_payment():
    """Handle offline payment"""
    payinfo = st.session_state.get("payinfo")
    if not payinfo:
        st.error("❌ No payment information available")
        return
    
    amount = payinfo.get("amount", 0)
    currency = payinfo.get("currency", "INR")
    
    st.success(f"💵 Please pay ₹{amount:.2f} {currency} offline to the print shop")
    st.info("Show this screen to the print shop operator as proof of your print job.")
    
    st.balloons()
    complete_payment_process()

def complete_payment_process():
    """Mark payment process as complete"""
    st.session_state.payinfo = None
    st.session_state.process_complete = True
    set_status("✅ Payment process completed!")

def cancel_payment():
    """Cancel the payment process"""
    st.session_state.payinfo = None
    set_status("❌ Payment cancelled by user")

# --------- Main UI ----------

# Sidebar for system information
with st.sidebar:
    st.title("📋 System Info")
    
    # Environment status
    with st.expander("🔧 Environment Status"):
        st.write(f"**Platform:** {platform.system()}")
        st.write(f"**Firestore:** {'✅ Connected' if FIRESTORE_OK else '❌ Not Available'}")
        st.write(f"**PDF Reader:** {'✅ Available' if PDF_READER_AVAILABLE else '❌ Not Available'}")
        st.write(f"**python-docx:** {'✅ Available' if PYTHON_DOCX_AVAILABLE else '❌ Not Available'}")
        st.write(f"**python-pptx:** {'✅ Available' if PYTHON_PPTX_AVAILABLE else '❌ Not Available'}")
        st.write(f"**QR Code:** {'✅ Available' if QR_AVAILABLE else '❌ Not Available'}")
    
    # Supported formats
    with st.expander("📄 Supported Formats"):
        st.write("**✅ Supported:**")
        st.write("• PDF (passthrough)")
        st.write("• Text: txt, md, rtf, html, csv")
        st.write("• Images: png, jpg, jpeg, bmp, tiff, webp")
        st.write("• Documents: docx, pptx")
        
        st.write("**❌ Not Supported:**")
        st.write("• Excel files (.xlsx, .xls)")
        st.write("• Word 97-2003 (.doc)")
        st.write("• PowerPoint 97-2003 (.ppt)")

# User Information Section
st.markdown("### 👤 User Information")
col1, col2 = st.columns([2, 1])

with col1:
    user_name = st.text_input(
        "Your Name (Optional)", 
        value=st.session_state.get("user_name", ""),
        placeholder="Enter your name for the print job"
    )
    st.session_state.user_name = user_name

with col2:
    st.text_input(
        "User ID", 
        value=st.session_state.user_id, 
        disabled=True,
        help="Unique identifier for your session"
    )

# File Upload Section
st.markdown("### 📂 File Upload")

if not FIRESTORE_OK:
    st.error(f"❌ Firestore not available: {FIRESTORE_ERR}")
    st.info("Please configure Firestore credentials in Streamlit secrets to enable file uploads.")
else:
    uploaded_files = st.file_uploader(
        "Choose files to print",
        accept_multiple_files=True,
        type=['pdf', 'txt', 'md', 'rtf', 'html', 'htm', 'png', 'jpg', 'jpeg', 'bmp', 'tiff', 'webp', 'docx', 'pptx', 'csv'],
        help="Select one or more files. Supported formats: PDF, Text, Images, DOCX, PPTX"
    )
    
    if uploaded_files:
        # Convert uploaded files
        with st.spinner("🔄 Converting files..."):
            converted_files = []
            conversion_results = []
            
            for uploaded_file in uploaded_files:
                try:
                    converted_file = FileConverter.convert_uploaded_file_to_pdf(uploaded_file)
                    if converted_file:
                        converted_files.append(converted_file)
                        conversion_results.append({
                            "filename": uploaded_file.name,
                            "status": "✅ Success",
                            "method": converted_file.conversion_method,
                            "pages": converted_file.pages
                        })
                    else:
                        conversion_results.append({
                            "filename": uploaded_file.name,
                            "status": "❌ Failed",
                            "method": "unknown",
                            "pages": 0
                        })
                except Exception as e:
                    logger.error(f"Conversion error for {uploaded_file.name}: {e}")
                    conversion_results.append({
                        "filename": uploaded_file.name,
                        "status": f"❌ Error: {str(e)[:50]}",
                        "method": "error",
                        "pages": 0
                    })
            
            st.session_state.converted_files = converted_files
        
        # Show conversion results
        if conversion_results:
            st.markdown("#### 📋 Conversion Results")
            
            for result in conversion_results:
                col1, col2, col3, col4 = st.columns([3, 1, 1, 1])
                
                with col1:
                    st.write(f"**{result['filename']}**")
                with col2:
                    st.write(result['status'])
                with col3:
                    st.write(f"Method: {result['method']}")
                with col4:
                    st.write(f"Pages: {result['pages']}")
        
        # File Preview Section
        if converted_files:
            st.markdown("#### 👀 File Preview")
            
            for i, cf in enumerate(converted_files):
                with st.expander(f"📄 {cf.pdf_name} ({cf.pages} pages)"):
                    col1, col2, col3 = st.columns([2, 1, 1])
                    
                    with col1:
                        st.write(f"**Original:** {cf.orig_name}")
                        st.write(f"**Converted:** {cf.pdf_name}")
                        st.write(f"**Method:** {cf.conversion_method}")
                        st.write(f"**Size:** {len(cf.pdf_bytes):,} bytes")
                    
                    with col2:
                        if st.button(f"👁️ Preview", key=f"preview_{i}"):
                            # Create inline PDF viewer
                            b64_pdf = base64.b64encode(cf.pdf_bytes).decode('utf-8')
                            pdf_display = f"""
                            <iframe src="data:application/pdf;base64,{b64_pdf}" 
                                    width="100%" height="600" type="application/pdf">
                            </iframe>
                            """
                            st.markdown(pdf_display, unsafe_allow_html=True)
                    
                    with col3:
                        st.download_button(
                            "💾 Download PDF",
                            data=cf.pdf_bytes,
                            file_name=cf.pdf_name,
                            mime="application/pdf",
                            key=f"download_{i}"
                        )
        
        # Print Job Settings
        if converted_files:
            st.markdown("#### ⚙️ Print Job Settings")
            
            col1, col2, col3 = st.columns(3)
            
            with col1:
                copies = st.number_input(
                    "Copies", 
                    min_value=1, 
                    max_value=20, 
                    value=1,
                    help="Number of copies for each file"
                )
            
            with col2:
                color_mode = st.selectbox(
                    "Color Mode",
                    options=["Auto", "Color", "Monochrome"],
                    help="Color printing mode"
                )
            
            with col3:
                paper_size = st.selectbox(
                    "Paper Size",
                    options=["A4", "A3", "Letter"],
                    help="Paper size for printing"
                )
            
            # Calculate total pages and estimated cost
            total_pages = sum(cf.pages * copies for cf in converted_files)
            pricing = st.session_state.pricing
            
            is_color = "color" in color_mode.lower()
            estimated_cost = calculate_amount(pricing, total_pages, 1, is_color, False)
            
            st.info(f"📊 **Total Pages:** {total_pages} | **Estimated Cost:** ₹{estimated_cost:.2f}")
            
            # Upload Button
            job_settings = {
                "copies": copies,
                "color_mode": color_mode,
                "paper_size": paper_size
            }
            
            if st.button("🚀 Send Files for Printing", type="primary", use_container_width=True):
                success = upload_files_to_firestore(converted_files, job_settings)
                if success:
                    st.success("✅ Files uploaded successfully!")

# Status Display
if st.session_state.get("status"):
    st.info(f"📊 **Status:** {st.session_state.status}")

# Payment Section
payinfo = st.session_state.get("payinfo")
if payinfo and not st.session_state.get("process_complete"):
    st.markdown("---")
    st.markdown("### 💳 Payment Required")
    
    # Payment details
    col1, col2 = st.columns(2)
    
    with col1:
        st.markdown("#### 📄 Order Details")
        st.write(f"**Files:** {payinfo.get('file_name', 'Multiple Files')}")
        if payinfo.get('total_files', 0) > 1:
            st.write(f"**Total Files:** {payinfo['total_files']}")
        st.write(f"**Pages:** {payinfo.get('pages', 'N/A')}")
        st.write(f"**Copies:** {payinfo.get('copies', 1)}")
        st.write(f"**Amount:** ₹{payinfo.get('amount', 0):.2f} {payinfo.get('currency', 'INR')}")
        
        if payinfo.get('is_estimate'):
            st.warning("⚠️ This is a local estimate. Official pricing may vary.")
    
    with col2:
        st.markdown("#### 💰 Payment Options")
        
        col_online, col_offline, col_cancel = st.columns(3)
        
        with col_online:
            if st.button("💳 Pay Online", type="primary", use_container_width=True):
                handle_online_payment()
        
        with col_offline:
            if st.button("💵 Pay at Shop", use_container_width=True):
                handle_offline_payment()
        
        with col_cancel:
            if st.button("❌ Cancel", use_container_width=True):
                cancel_payment()

# Process Complete Section
if st.session_state.get("process_complete"):
    st.markdown("---")
    st.success("🎉 **Print job submitted successfully!**")
    st.info("Your files have been sent to the print shop. Please proceed with payment and collect your prints.")
    
    if st.button("🔄 Start New Print Job", type="primary"):
        # Reset session state for new job
        st.session_state.converted_files = []
        st.session_state.payinfo = None
        st.session_state.process_complete = False
        st.session_state.status = ""
        st.session_state.user_id = str(uuid.uuid4())[:8]
        st.rerun()

# Footer
st.markdown("---")
st.markdown(
    "<div style='text-align: center; color: #666; padding: 1rem;'>"
    "🖨️ <strong>Autoprint</strong> - Streamlit Cloud Compatible Firestore Sender<br>"
    "<small>Optimized for cloud deployment with improved error handling and fallbacks</small>"
    "</div>", 
    unsafe_allow_html=True
)
